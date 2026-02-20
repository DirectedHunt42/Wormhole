import os
import sys
import shutil  # Added for ffmpeg/pandoc checks
import customtkinter as ctk
from tkinter import filedialog, messagebox
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from pypdf import PdfReader
import py7zr
import zipfile
import tarfile
import tempfile
import shutil
import threading
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
import openpyxl
import csv
import webbrowser
import subprocess
import re
import darkdetect
try:
    from striprtf.striprtf import rtf_to_text
    RTF_SUPPORT = True
except ImportError:
    RTF_SUPPORT = False
import ezodf
import requests
import json
try:
    import trimesh
    TRIMESH_SUPPORT = True
except ImportError:
    TRIMESH_SUPPORT = False

# Debug prints to diagnose environment
print("Python executable:", sys.executable)
print("Python version:", sys.version)
print("sys.path (search paths for imports):", sys.path)

if (darkdetect.theme() == "Light"):
    BG = "#f7f8ff"
    CARD = "#ffffff"
    CARD_HOVER = "#e7e9f5"
    ACCENT = "#3a63d9"
    ACCENT_DIM = "#2a4ba8"
    TEXT = "#1a1b25"
else:
    BG = "#0a0812"
    CARD = "#120f1e"
    CARD_HOVER = "#19172b"
    ACCENT = "#7aa3ff"
    ACCENT_DIM = "#4d6bbc"
    TEXT = "#e8e6f5"

def resource_path(relative_path):
    """Get absolute path to resource, works for dev and for PyInstaller onefile."""
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except AttributeError:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)
if (darkdetect.theme() == "Light"):
    WORMHOLE_IMAGE_PATH = resource_path(os.path.join("Icons", "wormhole_Transparent.png"))
else:
    WORMHOLE_IMAGE_PATH = resource_path(os.path.join("Icons", "wormhole_Transparent_Light.png"))
try:
    WORMHOLE_PIL_IMAGE = Image.open(WORMHOLE_IMAGE_PATH)
except Exception as e:
    print(f"Could not load wormhole image: {e}")
    WORMHOLE_PIL_IMAGE = Image.new("RGBA", (100, 100), (100, 100, 100, 255))

APP_ICON_PATH = resource_path(os.path.join("Icons", "Wormhole_Icon.ico"))

FONTS_DIR = resource_path("fonts")
FONT_FAMILY_REGULAR = "Pathway Extreme 36pt Regular"
FONT_FAMILY_SEMIBOLD = "Pathway Extreme 36pt SemiBold"
FONT_FAMILY_ITALIC = "Pathway Extreme 36pt Italic"
FONT_FAMILY_THIN = "Pathway Extreme 36pt Thin"
FONT_FAMILY_BLACK = "Pathway Extreme 36pt Black"
FONT_FILES = [
    "PathwayExtreme_36pt-Black.ttf",
    "PathwayExtreme_36pt-Italic.ttf",
    "PathwayExtreme_36pt-Regular.ttf",
    "PathwayExtreme_36pt-SemiBold.ttf",
    "PathwayExtreme_36pt-Thin.ttf"
]

VERSION = "1.2.5"
GITHUB_URL = "https://github.com/DirectedHunt42/Wormhole"

# Set up customtkinter
if (darkdetect.theme() == "Light"):
    ctk.set_appearance_mode("light")
    ctk.set_default_color_theme("blue")
else:
    ctk.set_appearance_mode("dark")
    ctk.set_default_color_theme("dark-blue")

# Load custom fonts
for font_file in FONT_FILES:
    font_path = os.path.join(FONTS_DIR, font_file)
    if os.path.exists(font_path):
        ctk.FontManager.load_font(font_path)
    else:
        print(f"Custom font file not found: {font_path}; falling back to default for this variant.")

has_pandoc = shutil.which("pandoc") is not None
has_ffmpeg = shutil.which("ffmpeg") is not None

formats = {
    'docs': {
        'extensions': [".txt", ".pdf", ".docx", ".html", ".md", ".odt"],
        'targets': ["TXT", "DOCX", "HTML", "MD", "ODT"],
    },
    'presentations': {
        'extensions': [".pptx", ".odp"],
        'targets': ["PPTX", "PDF", "TXT", "DOCX", "ODP"],
    },
    'images': {
        'extensions': [".jpg", ".jpeg", ".png", ".webp", ".avif", ".ico", ".bmp", ".gif", ".tiff"],
        'targets': ["PNG", "JPG", "WEBP", "AVIF", "ICO", "BMP", "GIF", "TIFF"],
    },
    'archive': {
        'extensions': [".zip", ".7z", ".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tbz2"],
        'targets': ["ZIP", "7Z", "TAR", "TGZ", "TBZ2"],
    },
    'spreadsheets': {
        'extensions': [".xlsx", ".csv", ".ods"],
        'targets': ["XLSX", "CSV", "ODS"],
    },
    '3d': {
        'extensions': [".obj", ".stl", ".ply", ".fbx", ".glb"],
        'targets': ["OBJ", "STL", "PLY", "FBX", "GLB"],
    },
    'media_audio': {
        'extensions': [".mp3", ".wav", ".ogg", ".flac", ".aac", ".m4a"],
        'targets': ["MP3", "WAV", "OGG", "FLAC", "AAC", "M4A"],
    },
    'media_video': {
        'extensions': [".mp4", ".avi", ".mkv", ".mov"],
        'targets': ["MP4", "AVI", "MKV", "MOV"] + [f"{a.upper()} (extract audio)" for a in ["mp3", "wav", "ogg", "flac", "aac", "m4a"]],
    },
}

if has_pandoc or RTF_SUPPORT:
    formats['docs']['extensions'].append('.rtf')
if has_pandoc:
    formats['docs']['targets'].append('RTF')
if not TRIMESH_SUPPORT:
    if '3d' in formats:
        del formats['3d']
if not has_ffmpeg:
    if 'media_audio' in formats:
        del formats['media_audio']
    if 'media_video' in formats:
        del formats['media_video']

def get_category(file_path):
    lfp = file_path.lower()
    audio_exts = ('.mp3', '.wav', '.ogg', '.flac', '.aac', '.m4a')
    video_exts = ('.mp4', '.avi', '.mkv', '.mov')
    if lfp.endswith(('.txt', '.pdf', '.docx', '.html', '.md', '.odt')) or (lfp.endswith('.rtf') and (has_pandoc or RTF_SUPPORT)):
        return 'docs'
    elif lfp.endswith(('.pptx', '.odp')):
        return 'presentations'
    elif lfp.endswith(('.jpg', '.jpeg', '.png', '.webp', '.avif', '.ico', '.bmp', '.gif', '.tiff')):
        return 'images'
    elif lfp.endswith(('.zip', '.7z', '.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2')):
        return 'archive'
    elif lfp.endswith(('.xlsx', '.csv', '.ods')):
        return 'spreadsheets'
    elif lfp.endswith(('.obj', '.stl', '.ply', '.fbx', '.glb')):
        return '3d'
    elif lfp.endswith(audio_exts):
        return 'media_audio'
    elif lfp.endswith(video_exts):
        return 'media_video'
    else:
        return None

def convert_docs(file_path, target):
    input_ext = os.path.splitext(file_path)[1].lower()[1:]
    new_file_path = os.path.splitext(file_path)[0] + '.' + target.lower()
    use_pandoc = has_pandoc and input_ext != "pdf" and target != "TXT" and target != "MD"
    text = ""
    if not use_pandoc:
        if input_ext in ["txt", "md"]:
            with open(file_path, 'r') as f:
                text = f.read()
        elif input_ext == "pdf":
            reader = PdfReader(file_path)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n'
        elif input_ext == "docx":
            doc = Document(file_path)
            text = '\n'.join([para.text for para in doc.paragraphs])
        elif input_ext == "html":
            with open(file_path, 'r') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                text = soup.get_text()
        elif input_ext == "odt":
            doc = ezodf.opendoc(file_path)
            text = '\n'.join(obj.text or '' for obj in doc.body if obj.kind == 'Paragraph')
        elif input_ext == "rtf":
            if RTF_SUPPORT:
                with open(file_path, 'r') as f:
                    rtf = f.read()
                text = rtf_to_text(rtf)
            else:
                raise ValueError("RTF input not supported without striprtf or Pandoc")
        else:
            raise ValueError("Unsupported input format")

        if target in ["TXT", "MD"]:
            with open(new_file_path, 'w') as f:
                f.write(text)
        elif target == "DOCX":
            doc = Document()
            for para_text in text.split('\n'):
                doc.add_paragraph(para_text)
            doc.save(new_file_path)
        elif target == "HTML":
            with open(new_file_path, 'w') as f:
                escaped_text = text.replace('<', '&lt;').replace('>', '&gt;')
                f.write(f"<html><body><pre>{escaped_text}</pre></body></html>")
        elif target == "ODT":
            doc = ezodf.newdoc(doctype='odt', filename=new_file_path)
            for para_text in text.split('\n'):
                doc.body.append(ezodf.Paragraph(para_text))
            doc.save()
        elif target == "RTF":
            raise ValueError("RTF output not supported without Pandoc")
        else:
            raise ValueError("Unsupported target format")
    else:
        subprocess.run(["pandoc", file_path, "-o", new_file_path], check=True)
    return new_file_path

def convert_presentations(file_path, target):
    input_ext = os.path.splitext(file_path)[1].lower()[1:]
    new_file_path = os.path.splitext(file_path)[0] + '.' + target.lower()
    text = ""
    if input_ext == "pptx":
        pres = Presentation(file_path)
        text = ''
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
                        text += '\n'
    elif input_ext == "odp":
        doc = ezodf.opendoc(file_path)
        text = '\n'.join(obj.text or '' for obj in doc.body if obj.kind == 'Paragraph')
    else:
        raise ValueError("Unsupported input format")

    if target == "TXT":
        with open(new_file_path, 'w') as f:
            f.write(text)
    elif target == "PDF":
        c = canvas.Canvas(new_file_path, pagesize=letter)
        width, height = letter
        y = height - 50  # Start from top with margin
        for line in text.splitlines():
            c.drawString(50, y, line.strip())
            y -= 15  # Line spacing
            if y < 50:  # Simple page break handling
                c.showPage()
                y = height - 50
        c.save()
    elif target == "DOCX":
        doc = Document()
        doc.add_paragraph(text)
        doc.save(new_file_path)
    elif target == "PPTX":
        pres = Presentation()
        slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(slide_layout)
        left = top = Inches(1)
        width = height = Inches(6)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text
        pres.save(new_file_path)
    elif target == "ODP":
        doc = ezodf.newdoc(doctype='odp', filename=new_file_path)
        doc.body.append(ezodf.Paragraph(text))
        doc.save()
    else:
        raise ValueError("Unsupported target format")
    return new_file_path

def convert_images(file_path, target):
    input_ext = os.path.splitext(file_path)[1].lower()[1:]
    if input_ext == "jpeg":
        input_ext = "jpg"
    new_file_path = os.path.splitext(file_path)[0] + '.' + target.lower()
    img = Image.open(file_path)
    if target == "JPG":
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')
        img.save(new_file_path, 'JPEG')
    elif target == "PNG":
        img.save(new_file_path, 'PNG')
    elif target == "WEBP":
        img.save(new_file_path, 'WEBP')
    elif target == "AVIF":
        img.save(new_file_path, 'AVIF')
    elif target == "ICO":
        sizes = [16, 24, 32, 40, 48, 64, 128, 256]
        selected_tuples = [(s, s) for s in sizes]
        img.save(new_file_path, format='ICO', sizes=selected_tuples, bitmap_format="bmp")
    elif target == "BMP":
        img.save(new_file_path, 'BMP')
    elif target == "GIF":
        img.save(new_file_path, 'GIF')
    elif target == "TIFF":
        img.save(new_file_path, 'TIFF')
    else:
        raise ValueError("Unsupported target format")
    return new_file_path

def get_archive_type(fp):
    lfp = fp.lower()
    if lfp.endswith('.zip'):
        return 'zip'
    if lfp.endswith('.7z'):
        return '7z'
    if lfp.endswith('.tar'):
        return 'tar'
    if lfp.endswith('.tar.gz') or lfp.endswith('.tgz'):
        return 'tgz'
    if lfp.endswith('.tar.bz2') or lfp.endswith('.tbz2'):
        return 'tbz2'
    raise ValueError("Unsupported archive type")

def convert_archive(file_path, target):
    input_type = get_archive_type(file_path)
    ext_map = {
        'ZIP': '.zip',
        '7Z': '.7z',
        'TAR': '.tar',
        'TGZ': '.tar.gz',
        'TBZ2': '.tar.bz2'
    }
    new_file_path = os.path.splitext(file_path)[0] + ext_map[target]
    temp_dir = tempfile.mkdtemp()
    try:
        # Extract
        if input_type == 'zip':
            with zipfile.ZipFile(file_path, 'r') as z:
                z.extractall(temp_dir)
        elif input_type == '7z':
            with py7zr.SevenZipFile(file_path, 'r') as z:
                z.extractall(temp_dir)
        elif input_type == 'tar':
            with tarfile.open(file_path, 'r') as t:
                t.extractall(temp_dir)
        elif input_type == 'tgz':
            with tarfile.open(file_path, 'r:gz') as t:
                t.extractall(temp_dir)
        elif input_type == 'tbz2':
            with tarfile.open(file_path, 'r:bz2') as t:
                t.extractall(temp_dir)
        
        # Create new archive
        if target == 'ZIP':
            with zipfile.ZipFile(new_file_path, 'w', zipfile.ZIP_DEFLATED) as z:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        z.write(os.path.join(root, file), os.path.relpath(os.path.join(root, file), temp_dir))
        elif target == '7Z':
            with py7zr.SevenZipFile(new_file_path, 'w') as z:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        z.write(os.path.join(root, file), os.path.relpath(os.path.join(root, file), temp_dir))
        elif target == 'TAR':
            with tarfile.open(new_file_path, 'w') as t:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        t.add(os.path.join(root, file), os.path.relpath(os.path.join(root, file), temp_dir))
        elif target == 'TGZ':
            with tarfile.open(new_file_path, 'w:gz') as t:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        t.add(os.path.join(root, file), os.path.relpath(os.path.join(root, file), temp_dir))
        elif target == 'TBZ2':
            with tarfile.open(new_file_path, 'w:bz2') as t:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        t.add(os.path.join(root, file), os.path.relpath(os.path.join(root, file), temp_dir))
    finally:
        shutil.rmtree(temp_dir)
    return new_file_path

def convert_spreadsheets(file_path, target):
    input_ext = os.path.splitext(file_path)[1].lower()[1:]
    new_file_path = os.path.splitext(file_path)[0] + '.' + target.lower()
    data = []
    if input_ext == "xlsx":
        wb = openpyxl.load_workbook(file_path)
        sheet = wb.active
        data = [[cell.value or '' for cell in row] for row in sheet.rows]
    elif input_ext == "csv":
        with open(file_path, 'r', newline='') as f:
            reader = csv.reader(f)
            data = list(reader)
    elif input_ext == "ods":
        doc = ezodf.opendoc(file_path)
        sheet = doc.sheets[0]
        data = [[cell.value or '' for cell in row] for row in sheet.rows()]
    else:
        raise ValueError("Unsupported input format")

    if target == "XLSX":
        wb = openpyxl.Workbook()
        sheet = wb.active
        for row in data:
            sheet.append(row)
        wb.save(new_file_path)
    elif target == "CSV":
        with open(new_file_path, 'w', newline='') as f:
            writer = csv.writer(f)
            writer.writerows(data)
    elif target == "ODS":
        doc = ezodf.newdoc(doctype='ods', filename=new_file_path)
        max_cols = max(len(row) for row in data) if data else 1
        sht = ezodf.Sheet('Sheet1', size=(len(data), max_cols))
        doc.sheets.append(sht)
        for r, row in enumerate(data):
            for c, val in enumerate(row):
                sht[r, c].set_value(val)
        doc.save()
    else:
        raise ValueError("Unsupported target format")
    return new_file_path

def convert_3d(file_path, target):
    if not TRIMESH_SUPPORT:
        raise ImportError("trimesh library not installed.")
    new_file_path = os.path.splitext(file_path)[0] + '.' + target.lower()
    mesh = trimesh.load(file_path)
    mesh.export(new_file_path)
    return new_file_path

def run_ffmpeg(cmd, progress_cb=None, duration=None):
    process = subprocess.Popen(cmd, stderr=subprocess.PIPE, universal_newlines=True)
    while True:
        line = process.stderr.readline().strip()
        if not line and process.poll() is not None:
            break
        if line and progress_cb and duration:
            time_match = re.search(r'time=(\d{2}):(\d{2}):(\d{2}\.\d{2})', line)
            if time_match:
                h, m, s = map(float, time_match.groups())
                time = h * 3600 + m * 60 + s
                progress_cb(time / duration)
    process.wait()
    if process.returncode != 0:
        raise RuntimeError(f"ffmpeg failed with code {process.returncode}")

def convert_media(file_path, target):
    if not has_ffmpeg:
        raise RuntimeError("ffmpeg not found in PATH.")
    is_extract = "(extract audio)" in target
    target_ext = target.split(" ")[0].lower()
    new_file_path = os.path.splitext(file_path)[0] + '.' + target_ext
    # Get duration using ffprobe
    duration = None
    try:
        ffprobe_cmd = ['ffprobe', '-v', 'error', '-show_entries', 'format=duration', '-of', 'default=noprint_wrappers=1:nokey=1', file_path]
        duration_str = subprocess.check_output(ffprobe_cmd).decode().strip()
        duration = float(duration_str)
    except Exception as e:
        print(f"Could not get duration: {e}")
    if is_extract:
        audio_cmd = ['ffmpeg', '-y', '-i', file_path, '-vn', new_file_path]
        muted_path = os.path.splitext(file_path)[0] + '_no_audio' + os.path.splitext(file_path)[1]
        video_cmd = ['ffmpeg', '-y', '-i', file_path, '-an', muted_path]
        if duration:
            def audio_prog(time):
                pass  # no progress in silent
            run_ffmpeg(audio_cmd, audio_prog, duration)
            def video_prog(time):
                pass
            run_ffmpeg(video_cmd, video_prog, duration)
        else:
            subprocess.check_call(audio_cmd)
            subprocess.check_call(video_cmd)
        return f"{new_file_path}, {muted_path}"
    else:
        cmd = ['ffmpeg', '-y', '-i', file_path, new_file_path]
        if duration:
            def conv_prog(time):
                pass
            run_ffmpeg(cmd, conv_prog, duration)
        else:
            subprocess.check_call(cmd)
        return new_file_path

def silent_convert(file_path, target):
    if not os.path.isfile(file_path):
        print("File not found")
        sys.exit(1)
    input_ext = os.path.splitext(file_path)[1].lower()[1:]
    if input_ext == 'jpeg':
        input_ext = 'jpg'
    cat = get_category(file_path)
    if not cat:
        print("Unsupported file type")
        sys.exit(1)
    if target.lower() == input_ext or (target.lower().startswith(input_ext) and "(extract audio)" in target):
        print("Input and output formats are the same")
        sys.exit(0)
    try:
        if cat == 'docs':
            new_fp = convert_docs(file_path, target)
        elif cat == 'presentations':
            new_fp = convert_presentations(file_path, target)
        elif cat == 'images':
            new_fp = convert_images(file_path, target)
        elif cat == 'archive':
            new_fp = convert_archive(file_path, target)
        elif cat == 'spreadsheets':
            new_fp = convert_spreadsheets(file_path, target)
        elif cat == '3d':
            new_fp = convert_3d(file_path, target)
        elif cat in ['media_audio', 'media_video']:
            new_fp = convert_media(file_path, target)
        # On success, stay silent (no print)
    except Exception as e:
        # Optionally show error popup (uncomment if wanted; otherwise silent fail)
        # messagebox.showerror("Conversion Failed", str(e))
        print(f"Conversion failed: {str(e)}")  # Fallback, but silent in --windowed mode
        sys.exit(1)

class WormholeApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("Wormhole File Converter")
        self.geometry("400x740")
        self.configure(fg_color=BG)
        # Center the main window
        self.update_idletasks()
        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()
        x = (screen_width // 2) - (400 // 2)
        y = (screen_height // 2) - (740 // 2) - 30
        self.geometry(f"400x740+{x}+{y}")
        self._build_ui()
        self.check_for_updates()
        if sys.platform.startswith('win'):
            pass  # Registration moved to --register
        if len(sys.argv) > 1:
            file = sys.argv[1]
            if os.path.isfile(file):
                cat = get_category(file)
                if cat:
                    open_func = getattr(self, f'open_{cat}_window')
                    open_func(preselected_file=file)

    def register_context_menu(self):
        try:
            import winreg
            if hasattr(sys, '_MEIPASS'):
                # Bundled mode
                exe_path = sys.executable
            else:
                # Script mode
                exe_path = f'"{sys.executable}" "{os.path.abspath(__file__)}"'
            icon_path = sys.executable if hasattr(sys, '_MEIPASS') else APP_ICON_PATH
            for cat, info in formats.items():
                for ext in info['extensions']:
                    main_key_path = rf"Software\Classes\SystemFileAssociations\{ext}\shell\WormholeConvert"
                    key = winreg.CreateKeyEx(winreg.HKEY_CURRENT_USER, main_key_path, 0, winreg.KEY_SET_VALUE)
                    winreg.SetValueEx(key, "MUIVerb", 0, winreg.REG_SZ, "Convert with Wormhole")
                    winreg.SetValueEx(key, "Icon", 0, winreg.REG_SZ, icon_path)
                    winreg.SetValueEx(key, "SubCommands", 0, winreg.REG_SZ, "")
                    # Create the shell subkey under main
                    shell_key_path = main_key_path + r"\shell"
                    winreg.CreateKeyEx(winreg.HKEY_CURRENT_USER, shell_key_path, 0, winreg.KEY_SET_VALUE)
                    for i, tgt in enumerate(sorted(info['targets'])):  # Sort for consistent order
                        clean_tgt = tgt.replace(' ', '')
                        sub_key_name = f"{i:02d}_{clean_tgt}"  # Prefix for ordering
                        sub_key_path = shell_key_path + rf"\{sub_key_name}"
                        sub_key = winreg.CreateKeyEx(winreg.HKEY_CURRENT_USER, sub_key_path, 0, winreg.KEY_SET_VALUE)
                        winreg.SetValueEx(sub_key, None, 0, winreg.REG_SZ, f"To {tgt}")
                        winreg.CloseKey(sub_key)
                        cmd_key_path = sub_key_path + r"\command"
                        cmd_key = winreg.CreateKeyEx(winreg.HKEY_CURRENT_USER, cmd_key_path, 0, winreg.KEY_SET_VALUE)
                        winreg.SetValueEx(cmd_key, None, 0, winreg.REG_SZ, f'{exe_path} "%1" "{tgt}"')
                        winreg.CloseKey(cmd_key)
                    winreg.CloseKey(key)
        except Exception as e:
            print(f"Failed to register context menu: {e}")

    def unregister_context_menu(self):
        try:
            import winreg
            for cat, info in formats.items():
                for ext in info['extensions']:
                    main_key_path = rf"Software\Classes\SystemFileAssociations\{ext}\shell\WormholeConvert"
                    self._delete_registry_key(winreg.HKEY_CURRENT_USER, main_key_path)
                    # Also clean up old flat structure if exists
                    for tgt in info['targets']:
                        clean_tgt = tgt.replace(' ', '')
                        old_sub_key_path = rf"Software\Classes\SystemFileAssociations\{ext}\shell\To{clean_tgt}"
                        self._delete_registry_key(winreg.HKEY_CURRENT_USER, old_sub_key_path)
        except Exception as e:
            print(f"Failed to unregister context menu: {e}")

    def _delete_registry_key(self, root, key_path):
        import winreg
        try:
            key = winreg.OpenKey(root, key_path, 0, winreg.KEY_ALL_ACCESS)
            num_subkeys, _, _ = winreg.QueryInfoKey(key)
            for i in range(num_subkeys):
                subkey_name = winreg.EnumKey(key, 0)
                self._delete_registry_key(root, key_path + "\\" + subkey_name)
            winreg.CloseKey(key)
            winreg.DeleteKey(root, key_path)
        except FileNotFoundError:
            pass  # Key already gone

    def _build_ui(self):
        if os.path.exists(APP_ICON_PATH):
            try:
                if sys.platform.startswith('win'):
                    import ctypes
                    ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID("wormhole.file.converter")
                self.iconbitmap(APP_ICON_PATH)
            except Exception as e:
                print(f"Could not set application icon: {e}")
                
        # Create a scrollable frame to hold all content
        scrollable_frame = ctk.CTkScrollableFrame(self, fg_color=BG)
        scrollable_frame.pack(fill='both', expand=True)

        # Custom label for instructions
        image = ctk.CTkImage(light_image=WORMHOLE_PIL_IMAGE, dark_image=WORMHOLE_PIL_IMAGE, size=(306, 204))
        img_label = ctk.CTkLabel(scrollable_frame, image=image, text="", fg_color=BG)
        img_label.pack(pady=10)

        label = ctk.CTkLabel(scrollable_frame, text="Select a file type category:", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
        label.pack(pady=20)

        # Buttons for each category (using semibold for buttons if desired; otherwise keep normal)
        btn_docs = ctk.CTkButton(scrollable_frame, text="Docs", command=self.open_docs_window, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=300, font=(FONT_FAMILY_SEMIBOLD, 20))
        btn_docs.pack(pady=5)

        btn_presentations = ctk.CTkButton(scrollable_frame, text="Presentations", command=self.open_presentations_window, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=300, font=(FONT_FAMILY_SEMIBOLD, 20))
        btn_presentations.pack(pady=5)

        btn_images = ctk.CTkButton(scrollable_frame, text="Images", command=self.open_images_window, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=300, font=(FONT_FAMILY_SEMIBOLD, 20))
        btn_images.pack(pady=5)

        btn_archive = ctk.CTkButton(scrollable_frame, text="Archive", command=self.open_archive_window, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=300, font=(FONT_FAMILY_SEMIBOLD, 20))
        btn_archive.pack(pady=5)

        btn_spreadsheets = ctk.CTkButton(scrollable_frame, text="Spreadsheets", command=self.open_spreadsheets_window, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=300, font=(FONT_FAMILY_SEMIBOLD, 20))
        btn_spreadsheets.pack(pady=5)

        btn_3d = ctk.CTkButton(scrollable_frame, text="3D Models", command=self.open_3d_window, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=300, font=(FONT_FAMILY_SEMIBOLD, 20))
        btn_3d.pack(pady=5)

        btn_media = ctk.CTkButton(scrollable_frame, text="Media", command=self.open_media_window, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=300, font=(FONT_FAMILY_SEMIBOLD, 20))
        btn_media.pack(pady=5)

        about_label = ctk.CTkLabel(scrollable_frame, text=f"Wormhole File Converter\nVersion {VERSION}\n© 2025-2026 Nova Foundry", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
        about_label.pack(pady=20)

        links_frame = ctk.CTkFrame(scrollable_frame, fg_color=BG)

        support_link = ctk.CTkLabel(
            links_frame, text="Support Nova Foundry",
            font=("Nunito", 12, "underline"), text_color=ACCENT,
            fg_color=BG, cursor="hand2"
        )
        support_link.pack(side="left", padx=10)

        official_link = ctk.CTkLabel(
            links_frame, text="Visit Official Website",
            font=("Nunito", 12, "underline"), text_color=ACCENT,
            fg_color=BG, cursor="hand2"
        )

        help_link = ctk.CTkLabel(links_frame, text="Help", font=("Nunito", 12, "underline"), text_color=ACCENT, fg_color=BG, cursor="hand2")
        help_link.pack(side="left", padx=10)

        official_link.pack(side="left", padx=10)

        links_frame.pack()


        def open_official_link(event):
            webbrowser.open_new("https://novafoundry.ca")
        def open_support_link(event):
            webbrowser.open_new("https://novafoundry.ca/support")
        def open_help_link(event):
            webbrowser.open_new("https://github.com/DirectedHunt42/Wormhole/wiki")
        support_link.bind("<Button-1>", open_support_link)
        official_link.bind("<Button-1>", open_official_link)
        help_link.bind("<Button-1>", open_help_link)

        # Hide the scrollbar permanently
        scrollable_frame._scrollbar.grid_forget()

    def check_for_updates(self):
        try:
            response = requests.get("https://api.github.com/repos/DirectedHunt42/Wormhole/releases/latest")
            if response.status_code == 200:
                data = response.json()
                latest_version = data['tag_name']
                if self.is_newer_version(latest_version, VERSION):
                    if messagebox.askyesno("Update Available", f"A new version {latest_version} is available. Do you want to download and install it?"):
                        self.download_and_install_update(data)
        except Exception:
            pass  # Fail silently if no internet or other issues

    def is_newer_version(self, latest, current):
        def parse(v):
            return tuple(int(x) for x in v.lstrip('v').split('.'))
        return parse(latest) > parse(current)

    def download_and_install_update(self, data):
        for asset in data['assets']:
            if asset['name'] == 'Wormhole_setup.exe':
                url = asset['browser_download_url']
                try:
                    response = requests.get(url, stream=True)
                    if response.status_code == 200:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.exe') as tmp:
                            for chunk in response.iter_content(chunk_size=1024):
                                tmp.write(chunk)
                            tmp_path = tmp.name
                        os.startfile(tmp_path)
                        sys.exit(0)
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to download or run update: {str(e)}")
                return
        messagebox.showerror("Error", "Wormhole_setup.exe not found in the latest release.")

# Functions to open subwindows for each category

def open_docs_window(master, preselected_file=None):
    docs_win = ctk.CTkToplevel(master)
    docs_win.title("Docs Conversions")
    docs_win.geometry("300x300")
    docs_win.configure(fg_color=BG)
    # Center the window
    docs_win.update_idletasks()
    screen_width = docs_win.winfo_screenwidth()
    screen_height = docs_win.winfo_screenheight()
    x = (screen_width // 2) - (300 // 2)
    y = (screen_height // 2) - (300 // 2)
    docs_win.geometry(f"300x300+{x}+{y}")
    # Set icon
    if os.path.exists(APP_ICON_PATH):
        try:
            docs_win.after(250, lambda: docs_win.iconbitmap(APP_ICON_PATH))
        except Exception as e:
            print(f"Could not set icon for docs window: {e}")
    # Make it transient and grab set to stay on top
    docs_win.transient(master)
    docs_win.grab_set()

    label = ctk.CTkLabel(docs_win, text="Docs Converter", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
    label.pack(pady=10)

    file_path_var = ctk.StringVar(value="")

    base_filetypes = "*.txt;*.pdf;*.docx;*.html;*.md;*.odt"
    if has_pandoc or RTF_SUPPORT:
        base_filetypes += ";*.rtf"
    filetypes = [("Docs files", base_filetypes)]

    def select_file():
        fp = filedialog.askopenfilename(title="Select Docs File", filetypes=filetypes)
        if fp:
            file_path_var.set(fp)
            file_label.configure(text=os.path.basename(fp))

    btn_select = ctk.CTkButton(docs_win, text="Select File", command=select_file, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_select.pack(pady=5)

    file_label = ctk.CTkLabel(docs_win, text="No file selected", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    file_label.pack(pady=5)

    target_var = ctk.StringVar(value="TXT")
    combo_values = ["TXT", "DOCX", "HTML", "MD", "ODT"]
    if has_pandoc:
        combo_values.append("RTF")
    combo = ctk.CTkComboBox(docs_win, values=combo_values, variable=target_var, font=(FONT_FAMILY_REGULAR, 10), width=250)
    combo.pack(pady=5)

    progress_bar = ctk.CTkProgressBar(docs_win, width=250, mode="indeterminate")
    # Initially not packed

    if preselected_file:
        file_path_var.set(preselected_file)
        file_label.configure(text=os.path.basename(preselected_file))

    def do_convert():
        fp = file_path_var.get()
        if not fp:
            messagebox.showerror("Error", "No file selected")
            return
        target = target_var.get().upper()
        input_ext = os.path.splitext(fp)[1].lower()[1:]
        if target.lower() == input_ext:
            messagebox.showwarning("Warning", "Input and output formats are the same")
            return

        def conversion_thread():
            try:
                new_file_path = convert_docs(fp, target)
                docs_win.after(0, lambda: messagebox.showinfo("Success", f"File converted to: {new_file_path}"))
            except FileNotFoundError:
                docs_win.after(0, lambda: messagebox.showerror("Error", "Pandoc not found. Please install Pandoc for full formatting support."))
            except Exception as e:
                docs_win.after(0, lambda e=e: messagebox.showerror("Error", f"Conversion failed: {str(e)}"))
            finally:
                docs_win.after(0, progress_bar.stop)
                docs_win.after(0, progress_bar.pack_forget)
                docs_win.after(0, lambda: btn_convert.configure(state="normal"))

        progress_bar.pack(pady=5)
        progress_bar.start()
        btn_convert.configure(state="disabled")
        thread = threading.Thread(target=conversion_thread)
        thread.start()

    btn_convert = ctk.CTkButton(docs_win, text="Convert", command=do_convert, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_convert.pack(pady=5)

def open_presentations_window(master, preselected_file=None):
    pres_win = ctk.CTkToplevel(master)
    pres_win.title("Presentations Conversions")
    pres_win.geometry("300x300")
    pres_win.configure(fg_color=BG)
    # Center the window
    pres_win.update_idletasks()
    screen_width = pres_win.winfo_screenwidth()
    screen_height = pres_win.winfo_screenheight()
    x = (screen_width // 2) - (300 // 2)
    y = (screen_height // 2) - (300 // 2)
    pres_win.geometry(f"300x300+{x}+{y}")
    # Set icon
    if os.path.exists(APP_ICON_PATH):
        try:
            pres_win.after(250, lambda: pres_win.iconbitmap(APP_ICON_PATH))
        except Exception as e:
            print(f"Could not set icon for presentations window: {e}")
    # Make it transient and grab set to stay on top
    pres_win.transient(master)
    pres_win.grab_set()

    label = ctk.CTkLabel(pres_win, text="Presentations Converter", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
    label.pack(pady=10)

    file_path_var = ctk.StringVar(value="")

    def select_file():
        fp = filedialog.askopenfilename(title="Select Presentation File", filetypes=[("Presentation files", "*.pptx;*.odp")])
        if fp:
            file_path_var.set(fp)
            file_label.configure(text=os.path.basename(fp))

    btn_select = ctk.CTkButton(pres_win, text="Select File", command=select_file, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_select.pack(pady=5)

    file_label = ctk.CTkLabel(pres_win, text="No file selected", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    file_label.pack(pady=5)

    target_var = ctk.StringVar(value="PDF")
    combo = ctk.CTkComboBox(pres_win, values=["PPTX", "PDF", "TXT", "DOCX", "ODP"], variable=target_var, font=(FONT_FAMILY_REGULAR, 10), width=250)
    combo.pack(pady=5)

    progress_bar = ctk.CTkProgressBar(pres_win, width=250, mode="indeterminate")
    # Initially not packed

    if preselected_file:
        file_path_var.set(preselected_file)
        file_label.configure(text=os.path.basename(preselected_file))

    def do_convert():
        fp = file_path_var.get()
        if not fp:
            messagebox.showerror("Error", "No file selected")
            return
        target = target_var.get().upper()
        input_ext = os.path.splitext(fp)[1].lower()[1:]
        if target.lower() == input_ext:
            messagebox.showwarning("Warning", "Input and output formats are the same")
            return

        def conversion_thread():
            try:
                new_file_path = convert_presentations(fp, target)
                pres_win.after(0, lambda: messagebox.showinfo("Success", f"File converted to: {new_file_path}"))
            except Exception as e:
                pres_win.after(0, lambda e=e: messagebox.showerror("Error", f"Conversion failed: {str(e)}"))
            finally:
                pres_win.after(0, progress_bar.stop)
                pres_win.after(0, progress_bar.pack_forget)
                pres_win.after(0, lambda: btn_convert.configure(state="normal"))

        progress_bar.pack(pady=5)
        progress_bar.start()
        btn_convert.configure(state="disabled")
        thread = threading.Thread(target=conversion_thread)
        thread.start()

    btn_convert = ctk.CTkButton(pres_win, text="Convert", command=do_convert, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_convert.pack(pady=5)

def open_images_window(master, preselected_file=None):
    img_win = ctk.CTkToplevel(master)
    img_win.title("Images Conversions")
    img_win.geometry("300x450")
    img_win.configure(fg_color=BG)
    # Center the window
    img_win.update_idletasks()
    screen_width = img_win.winfo_screenwidth()
    screen_height = img_win.winfo_screenheight()
    x = (screen_width // 2) - (300 // 2)
    y = (screen_height // 2) - (450 // 2)
    img_win.geometry(f"300x450+{x}+{y}")
    # Set icon
    if os.path.exists(APP_ICON_PATH):
        try:
            img_win.after(250, lambda: img_win.iconbitmap(APP_ICON_PATH))
        except Exception as e:
            print(f"Could not set icon for images window: {e}")
    # Make it transient and grab set to stay on top
    img_win.transient(master)
    img_win.grab_set()

    label = ctk.CTkLabel(img_win, text="Images Converter", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
    label.pack(pady=10)

    file_path_var = ctk.StringVar(value="")

    def select_file():
        fp = filedialog.askopenfilename(title="Select Image File", filetypes=[("Image files", "*.jpg;*.jpeg;*.png;*.webp;*.avif;*.ico;*.bmp;*.gif;*.tiff")])
        if fp:
            file_path_var.set(fp)
            file_label.configure(text=os.path.basename(fp))

    btn_select = ctk.CTkButton(img_win, text="Select File", command=select_file, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_select.pack(pady=5)

    file_label = ctk.CTkLabel(img_win, text="No file selected", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    file_label.pack(pady=5)

    target_var = ctk.StringVar(value="PNG")
    combo = ctk.CTkComboBox(img_win, values=["PNG", "JPG", "WEBP", "AVIF", "ICO", "BMP", "GIF", "TIFF"], variable=target_var, font=(FONT_FAMILY_REGULAR, 10), width=250)
    combo.pack(pady=5)

    ico_frame = ctk.CTkFrame(img_win, fg_color=BG)
    ico_label = ctk.CTkLabel(ico_frame, text="Select ICO sizes:", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    ico_label.pack(pady=10)
    sizes = [16, 24, 32, 40, 48, 64, 128, 256]
    check_vars = [ctk.BooleanVar(value=True) for _ in sizes]  # Default checked
    for i, s in enumerate(sizes):
        cb = ctk.CTkCheckBox(ico_frame, text=f"{s}x{s}", variable=check_vars[i])
        cb.pack(anchor="w")

    def update_ico_frame(event=None):
        if target_var.get() == "ICO":
            ico_frame.pack(pady=5)
            img_win.geometry("300x450")
        else:
            ico_frame.pack_forget()
            img_win.geometry("300x300")
        img_win.update_idletasks()

    combo.configure(command=lambda choice: update_ico_frame())

    progress_bar = ctk.CTkProgressBar(img_win, width=250, mode="indeterminate")
    # Initially not packed

    if preselected_file:
        file_path_var.set(preselected_file)
        file_label.configure(text=os.path.basename(preselected_file))

    def do_convert():
        fp = file_path_var.get()
        if not fp:
            messagebox.showerror("Error", "No file selected")
            return
        target = target_var.get().upper()
        input_ext = os.path.splitext(fp)[1].lower()[1:]
        if input_ext == "jpeg":
            input_ext = "jpg"
        if target.lower() == input_ext or (target == "JPG" and input_ext in ["jpg", "jpeg"]):
            messagebox.showwarning("Warning", "Input and output formats are the same")
            return

        def conversion_thread():
            try:
                if target == "ICO":
                    selected_sizes = [sizes[i] for i, v in enumerate(check_vars) if v.get()]
                    if not selected_sizes:
                        img_win.after(0, lambda: messagebox.showerror("Error", "Select at least one size for ICO"))
                        return
                new_file_path = convert_images(fp, target)
                img_win.after(0, lambda: messagebox.showinfo("Success", f"File converted to: {new_file_path}"))
            except Exception as e:
                img_win.after(0, lambda e=e: messagebox.showerror("Error", f"Conversion failed: {str(e)}"))
            finally:
                img_win.after(0, progress_bar.stop)
                img_win.after(0, progress_bar.pack_forget)
                img_win.after(0, lambda: btn_convert.configure(state="normal"))

        progress_bar.pack(pady=5)
        progress_bar.start()
        btn_convert.configure(state="disabled")
        thread = threading.Thread(target=conversion_thread)
        thread.start()

    btn_convert = ctk.CTkButton(img_win, text="Convert", command=do_convert, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_convert.pack(pady=5)

    # Initially hide ico_frame if not ICO
    update_ico_frame()

def open_archive_window(master, preselected_file=None):
    arch_win = ctk.CTkToplevel(master)
    arch_win.title("Archive Conversions")
    arch_win.geometry("300x300")
    arch_win.configure(fg_color=BG)
    # Center the window
    arch_win.update_idletasks()
    screen_width = arch_win.winfo_screenwidth()
    screen_height = arch_win.winfo_screenheight()
    x = (screen_width // 2) - (300 // 2)
    y = (screen_height // 2) - (300 // 2)
    arch_win.geometry(f"300x300+{x}+{y}")
    # Set icon
    if os.path.exists(APP_ICON_PATH):
        try:
            arch_win.after(250, lambda: arch_win.iconbitmap(APP_ICON_PATH))
        except Exception as e:
            print(f"Could not set icon for archive window: {e}")
    # Make it transient and grab set to stay on top
    arch_win.transient(master)
    arch_win.grab_set()

    label = ctk.CTkLabel(arch_win, text="Archive Converter", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
    label.pack(pady=10)

    file_path_var = ctk.StringVar(value="")

    def select_file():
        fp = filedialog.askopenfilename(title="Select Archive File", filetypes=[("Archive files", "*.zip;*.7z;*.tar;*.tar.gz;*.tgz;*.tar.bz2;*.tbz2")])
        if fp:
            file_path_var.set(fp)
            file_label.configure(text=os.path.basename(fp))

    btn_select = ctk.CTkButton(arch_win, text="Select File", command=select_file, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_select.pack(pady=5)

    file_label = ctk.CTkLabel(arch_win, text="No file selected", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    file_label.pack(pady=5)

    target_var = ctk.StringVar(value="ZIP")
    combo = ctk.CTkComboBox(arch_win, values=["ZIP", "7Z", "TAR", "TGZ", "TBZ2"], variable=target_var, font=(FONT_FAMILY_REGULAR, 10), width=250)
    combo.pack(pady=5)

    progress_bar = ctk.CTkProgressBar(arch_win, width=250, mode="indeterminate")
    # Initially not packed

    if preselected_file:
        file_path_var.set(preselected_file)
        file_label.configure(text=os.path.basename(preselected_file))

    def do_convert():
        fp = file_path_var.get()
        if not fp:
            messagebox.showerror("Error", "No file selected")
            return
        target = target_var.get().upper()
        try:
            input_type = get_archive_type(fp)
        except ValueError:
            messagebox.showerror("Error", "Unsupported input format")
            return
        if target.lower() == input_type:
            messagebox.showwarning("Warning", "Input and output formats are the same")
            return

        def conversion_thread():
            try:
                new_file_path = convert_archive(fp, target)
                arch_win.after(0, lambda: messagebox.showinfo("Success", f"File converted to: {new_file_path}"))
            except Exception as e:
                arch_win.after(0, lambda e=e: messagebox.showerror("Error", f"Conversion failed: {str(e)}"))
            finally:
                arch_win.after(0, progress_bar.stop)
                arch_win.after(0, progress_bar.pack_forget)
                arch_win.after(0, lambda: btn_convert.configure(state="normal"))

        progress_bar.pack(pady=5)
        progress_bar.start()
        btn_convert.configure(state="disabled")
        thread = threading.Thread(target=conversion_thread)
        thread.start()

    btn_convert = ctk.CTkButton(arch_win, text="Convert", command=do_convert, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_convert.pack(pady=5)

def open_spreadsheets_window(master, preselected_file=None):
    spreadsheets_win = ctk.CTkToplevel(master)
    spreadsheets_win.title("Spreadsheets Conversions")
    spreadsheets_win.geometry("300x300")
    spreadsheets_win.configure(fg_color=BG)
    # Center the window
    spreadsheets_win.update_idletasks()
    screen_width = spreadsheets_win.winfo_screenwidth()
    screen_height = spreadsheets_win.winfo_screenheight()
    x = (screen_width // 2) - (300 // 2)
    y = (screen_height // 2) - (300 // 2)
    spreadsheets_win.geometry(f"300x300+{x}+{y}")
    # Set icon
    if os.path.exists(APP_ICON_PATH):
        try:
            spreadsheets_win.after(250, lambda: spreadsheets_win.iconbitmap(APP_ICON_PATH))
        except Exception as e:
            print(f"Could not set icon for spreadsheets window: {e}")
    # Make it transient and grab set to stay on top
    spreadsheets_win.transient(master)
    spreadsheets_win.grab_set()

    label = ctk.CTkLabel(spreadsheets_win, text="Spreadsheets Converter", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
    label.pack(pady=10)

    file_path_var = ctk.StringVar(value="")

    def select_file():
        fp = filedialog.askopenfilename(title="Select Spreadsheet File", filetypes=[("Spreadsheet files", "*.xlsx;*.csv;*.ods")])
        if fp:
            file_path_var.set(fp)
            file_label.configure(text=os.path.basename(fp))

    btn_select = ctk.CTkButton(spreadsheets_win, text="Select File", command=select_file, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_select.pack(pady=5)

    file_label = ctk.CTkLabel(spreadsheets_win, text="No file selected", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    file_label.pack(pady=5)

    target_var = ctk.StringVar(value="XLSX")
    combo = ctk.CTkComboBox(spreadsheets_win, values=["XLSX", "CSV", "ODS"], variable=target_var, font=(FONT_FAMILY_REGULAR, 10), width=250)
    combo.pack(pady=5)

    progress_bar = ctk.CTkProgressBar(spreadsheets_win, width=250, mode="indeterminate")
    # Initially not packed

    if preselected_file:
        file_path_var.set(preselected_file)
        file_label.configure(text=os.path.basename(preselected_file))

    def do_convert():
        fp = file_path_var.get()
        if not fp:
            messagebox.showerror("Error", "No file selected")
            return
        target = target_var.get().upper()
        input_ext = os.path.splitext(fp)[1].lower()[1:]
        if target.lower() == input_ext:
            messagebox.showwarning("Warning", "Input and output formats are the same")
            return

        def conversion_thread():
            try:
                new_file_path = convert_spreadsheets(fp, target)
                spreadsheets_win.after(0, lambda: messagebox.showinfo("Success", f"File converted to: {new_file_path}"))
            except Exception as e:
                spreadsheets_win.after(0, lambda e=e: messagebox.showerror("Error", f"Conversion failed: {str(e)}"))
            finally:
                spreadsheets_win.after(0, progress_bar.stop)
                spreadsheets_win.after(0, progress_bar.pack_forget)
                spreadsheets_win.after(0, lambda: btn_convert.configure(state="normal"))

        progress_bar.pack(pady=5)
        progress_bar.start()
        btn_convert.configure(state="disabled")
        thread = threading.Thread(target=conversion_thread)
        thread.start()

    btn_convert = ctk.CTkButton(spreadsheets_win, text="Convert", command=do_convert, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_convert.pack(pady=5)

def open_3d_window(master, preselected_file=None):
    if not TRIMESH_SUPPORT:
        messagebox.showerror("Error", "trimesh library not installed. Please install trimesh to enable 3D file support.")
        return

    threed_win = ctk.CTkToplevel(master)
    threed_win.title("3D Model Conversions")
    threed_win.geometry("300x300")
    threed_win.configure(fg_color=BG)
    # Center the window
    threed_win.update_idletasks()
    screen_width = threed_win.winfo_screenwidth()
    screen_height = threed_win.winfo_screenheight()
    x = (screen_width // 2) - (300 // 2)
    y = (screen_height // 2) - (300 // 2)
    threed_win.geometry(f"300x300+{x}+{y}")
    # Set icon
    if os.path.exists(APP_ICON_PATH):
        try:
            threed_win.after(250, lambda: threed_win.iconbitmap(APP_ICON_PATH))
        except Exception as e:
            print(f"Could not set icon for 3D window: {e}")
    # Make it transient and grab set to stay on top
    threed_win.transient(master)
    threed_win.grab_set()

    label = ctk.CTkLabel(threed_win, text="3D Model Converter", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
    label.pack(pady=10)

    file_path_var = ctk.StringVar(value="")

    def select_file():
        fp = filedialog.askopenfilename(title="Select 3D File", filetypes=[("3D files", "*.obj;*.stl;*.ply;*.fbx;*.glb")])
        if fp:
            file_path_var.set(fp)
            file_label.configure(text=os.path.basename(fp))

    btn_select = ctk.CTkButton(threed_win, text="Select File", command=select_file, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_select.pack(pady=5)

    file_label = ctk.CTkLabel(threed_win, text="No file selected", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    file_label.pack(pady=5)

    target_var = ctk.StringVar(value="OBJ")
    combo = ctk.CTkComboBox(threed_win, values=["OBJ", "STL", "PLY", "FBX", "GLB"], variable=target_var, font=(FONT_FAMILY_REGULAR, 10), width=250)
    combo.pack(pady=5)

    progress_bar = ctk.CTkProgressBar(threed_win, width=250, mode="indeterminate")
    # Initially not packed

    if preselected_file:
        file_path_var.set(preselected_file)
        file_label.configure(text=os.path.basename(preselected_file))

    def do_convert():
        fp = file_path_var.get()
        if not fp:
            messagebox.showerror("Error", "No file selected")
            return
        target = target_var.get().upper()
        input_ext = os.path.splitext(fp)[1].lower()[1:]
        if target.lower() == input_ext:
            messagebox.showwarning("Warning", "Input and output formats are the same")
            return

        def conversion_thread():
            try:
                new_file_path = convert_3d(fp, target)
                threed_win.after(0, lambda: messagebox.showinfo("Success", f"File converted to: {new_file_path}"))
            except Exception as e:
                threed_win.after(0, lambda e=e: messagebox.showerror("Error", f"Conversion failed: {str(e)}"))
            finally:
                threed_win.after(0, progress_bar.stop)
                threed_win.after(0, progress_bar.pack_forget)
                threed_win.after(0, lambda: btn_convert.configure(state="normal"))

        progress_bar.pack(pady=5)
        progress_bar.start()
        btn_convert.configure(state="disabled")
        thread = threading.Thread(target=conversion_thread)
        thread.start()

    btn_convert = ctk.CTkButton(threed_win, text="Convert", command=do_convert, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_convert.pack(pady=5)

def open_media_window(master, preselected_file=None):
    media_win = ctk.CTkToplevel(master)
    media_win.title("Media Conversions")
    media_win.geometry("300x300")
    media_win.configure(fg_color=BG)
    # Center the window
    media_win.update_idletasks()
    screen_width = media_win.winfo_screenwidth()
    screen_height = media_win.winfo_screenheight()
    x = (screen_width // 2) - (300 // 2)
    y = (screen_height // 2) - (300 // 2)
    media_win.geometry(f"300x300+{x}+{y}")
    # Set icon
    if os.path.exists(APP_ICON_PATH):
        try:
            media_win.after(250, lambda: media_win.iconbitmap(APP_ICON_PATH))
        except Exception as e:
            print(f"Could not set icon for media window: {e}")
    # Make it transient and grab set to stay on top
    media_win.transient(master)
    media_win.grab_set()

    label = ctk.CTkLabel(media_win, text="Media Converter", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 12))
    label.pack(pady=10)

    file_path_var = ctk.StringVar(value="")

    audio_formats = ["mp3", "wav", "ogg", "flac", "aac", "m4a"]
    video_formats = ["mp4", "avi", "mkv", "mov"]
    media_filetypes = [("Media files", "*." + ";*.".join(audio_formats + video_formats))]

    def select_file():
        fp = filedialog.askopenfilename(title="Select Media File", filetypes=media_filetypes)
        if fp:
            input_ext = os.path.splitext(fp)[1].lower()[1:]
            if input_ext in audio_formats:
                combo.configure(values=[fmt.upper() for fmt in audio_formats])
                target_var.set("MP3")
            elif input_ext in video_formats:
                video_values = [fmt.upper() for fmt in video_formats]
                audio_values = [fmt.upper() + " (extract audio)" for fmt in audio_formats]
                combo.configure(values=video_values + audio_values)
                target_var.set("MP4")
            else:
                messagebox.showerror("Error", "Unsupported media format")
                return
            file_path_var.set(fp)
            file_label.configure(text=os.path.basename(fp))

    btn_select = ctk.CTkButton(media_win, text="Select File", command=select_file, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_select.pack(pady=5)

    file_label = ctk.CTkLabel(media_win, text="No file selected", fg_color=BG, text_color=TEXT, font=(FONT_FAMILY_REGULAR, 10))
    file_label.pack(pady=5)

    target_var = ctk.StringVar(value="")
    combo = ctk.CTkComboBox(media_win, values=[], variable=target_var, font=(FONT_FAMILY_REGULAR, 10), width=250)
    combo.pack(pady=5)

    progress_bar = ctk.CTkProgressBar(media_win, width=250, mode="indeterminate")
    # Initially not packed

    if preselected_file:
        input_ext = os.path.splitext(preselected_file)[1].lower()[1:]
        if input_ext in audio_formats:
            combo.configure(values=[fmt.upper() for fmt in audio_formats])
            target_var.set("MP3")
        elif input_ext in video_formats:
            video_values = [fmt.upper() for fmt in video_formats]
            audio_values = [fmt.upper() + " (extract audio)" for fmt in audio_formats]
            combo.configure(values=video_values + audio_values)
            target_var.set("MP4")
        file_path_var.set(preselected_file)
        file_label.configure(text=os.path.basename(preselected_file))

    def do_convert():
        fp = file_path_var.get()
        if not fp:
            messagebox.showerror("Error", "No file selected")
            return
        target = target_var.get()
        input_ext = os.path.splitext(fp)[1].lower()[1:]
        is_extract = "(extract audio)" in target
        target_ext = target.split(" ")[0].lower()
        if target_ext == input_ext:
            messagebox.showwarning("Warning", "Input and output formats are the same")
            return

        def conversion_thread():
            try:
                new_file_path = convert_media(fp, target)
                media_win.after(0, lambda: messagebox.showinfo("Success", f"File converted to: {new_file_path}"))
            except Exception as e:
                media_win.after(0, lambda e=e: messagebox.showerror("Error", f"Conversion failed: {str(e)}"))
            finally:
                media_win.after(0, progress_bar.stop)
                media_win.after(0, progress_bar.pack_forget)
                media_win.after(0, lambda: btn_convert.configure(state="normal"))

        progress_bar.pack(pady=5)
        progress_bar.start()
        btn_convert.configure(state="disabled")
        thread = threading.Thread(target=conversion_thread)
        thread.start()

    btn_convert = ctk.CTkButton(media_win, text="Convert", command=do_convert, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, corner_radius=20, width=250, font=(FONT_FAMILY_SEMIBOLD, 10))
    btn_convert.pack(pady=5)

# Extend the app class with open methods
class WormholeApp(WormholeApp):
    def open_docs_window(self, preselected_file=None):
        open_docs_window(self, preselected_file)

    def open_presentations_window(self, preselected_file=None):
        open_presentations_window(self, preselected_file)

    def open_images_window(self, preselected_file=None):
        open_images_window(self, preselected_file)

    def open_archive_window(self, preselected_file=None):
        open_archive_window(self, preselected_file)

    def open_spreadsheets_window(self, preselected_file=None):
        open_spreadsheets_window(self, preselected_file)

    def open_3d_window(self, preselected_file=None):
        open_3d_window(self, preselected_file)

    def open_media_window(self, preselected_file=None):
        open_media_window(self, preselected_file)

if __name__ == "__main__":
    if len(sys.argv) > 1:
        if sys.argv[1] == "--register":
            app = WormholeApp()  # Init to access methods
            app.register_context_menu()
            sys.exit(0)
        elif sys.argv[1] == "--unregister":
            app = WormholeApp()
            app.unregister_context_menu()
            sys.exit(0)
        elif len(sys.argv) == 3:
            silent_convert(sys.argv[1], sys.argv[2])
            sys.exit(0)
    app = WormholeApp()
    app.mainloop()